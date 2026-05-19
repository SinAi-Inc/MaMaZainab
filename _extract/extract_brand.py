"""Extract brand assets: render PDF brand book to PNG, sample dominant HEX colors,
parse PPTX text, and extract MP4 keyframes."""
import os, json, collections
from pathlib import Path
import fitz  # PyMuPDF
from PIL import Image
import cv2
from pptx import Presentation

ROOT = Path(r"F:\H.Q\SinAI Inc\R&D Docs\Mama Zainab")
OUT = Path(r"d:\AI\MaMaZainab\_extract")
OUT.mkdir(parents=True, exist_ok=True)
(OUT / "pdf_pages").mkdir(exist_ok=True)
(OUT / "video_frames").mkdir(exist_ok=True)

PDF = ROOT / "BrandGuidelines" / "Mama Zainab Brand_guideline.pdf"
PDF2 = ROOT / "Archive" / "Phase 01-20260430T230237Z-3-001" / "Phase 01" / "Step 01 (Brand identity)" / "PDF" / "Mama Zainab Final version.pdf"
PPTX = ROOT / "BrandGuidelines" / "mamazienab.pptx"
MP4 = ROOT / "Archive" / "WhatsApp Video 2026-05-02 at 1.58.24 PM.mp4"

results = {"pdfs": [], "pptx": {}, "video": {}, "palette": {}}

# ---- PDFs ----
def render_pdf(pdf_path, tag):
    if not pdf_path.exists():
        return
    doc = fitz.open(pdf_path)
    info = {"file": str(pdf_path), "pages": len(doc), "rendered": []}
    for i, page in enumerate(doc):
        pix = page.get_pixmap(dpi=150)
        out = OUT / "pdf_pages" / f"{tag}_p{i+1:02d}.png"
        pix.save(out)
        info["rendered"].append(str(out))
        # also extract any embedded text
        text = page.get_text().strip()
        if text:
            (OUT / "pdf_pages" / f"{tag}_p{i+1:02d}.txt").write_text(text, encoding="utf-8")
    doc.close()
    results["pdfs"].append(info)

render_pdf(PDF, "guideline")
render_pdf(PDF2, "final")

# ---- PPTX ----
if PPTX.exists():
    prs = Presentation(PPTX)
    slides = []
    for idx, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    line = "".join(r.text for r in p.runs).strip()
                    if line:
                        texts.append(line)
        slides.append({"slide": idx + 1, "text": texts})
    results["pptx"] = {"file": str(PPTX), "slide_count": len(slides), "slides": slides}
    (OUT / "pptx_text.json").write_text(json.dumps(results["pptx"], ensure_ascii=False, indent=2), encoding="utf-8")

# ---- Video keyframes ----
if MP4.exists():
    cap = cv2.VideoCapture(str(MP4))
    fps = cap.get(cv2.CAP_PROP_FPS) or 30
    total = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
    duration = total / fps if fps else 0
    n_frames = 12
    saved = []
    for i in range(n_frames):
        target = int((i / max(n_frames - 1, 1)) * (total - 1))
        cap.set(cv2.CAP_PROP_POS_FRAMES, target)
        ok, frame = cap.read()
        if ok:
            out = OUT / "video_frames" / f"frame_{i:02d}.jpg"
            cv2.imwrite(str(out), frame, [int(cv2.IMWRITE_JPEG_QUALITY), 85])
            saved.append(str(out))
    cap.release()
    results["video"] = {"file": str(MP4), "fps": fps, "duration_s": round(duration, 2), "frames": saved}

# ---- Palette extraction (from rendered PDF pages) ----
def quantize_palette(img_path, k=8):
    img = Image.open(img_path).convert("RGB").resize((300, 300))
    q = img.quantize(colors=k, method=Image.Quantize.MEDIANCUT)
    pal = q.getpalette()[: k * 3]
    counts = collections.Counter(q.getdata())
    out = []
    for idx, cnt in counts.most_common():
        r, g, b = pal[idx * 3 : idx * 3 + 3]
        out.append({"hex": f"#{r:02X}{g:02X}{b:02X}", "rgb": [r, g, b], "weight": cnt})
    return out

palette_aggregate = collections.Counter()
per_page = {}
for png in sorted((OUT / "pdf_pages").glob("*.png")):
    pal = quantize_palette(png, k=8)
    per_page[png.name] = pal
    for c in pal:
        palette_aggregate[c["hex"]] += c["weight"]

results["palette"] = {
    "per_page": per_page,
    "top_overall": [{"hex": h, "weight": w} for h, w in palette_aggregate.most_common(20)],
}

(OUT / "extraction_report.json").write_text(json.dumps(results, ensure_ascii=False, indent=2), encoding="utf-8")
print("DONE")
print("PDF pages:", sum(len(p['rendered']) for p in results['pdfs']))
print("PPTX slides:", results['pptx'].get('slide_count', 0))
print("Video frames:", len(results['video'].get('frames', [])))
print("Top 12 colors:", palette_aggregate.most_common(12))
